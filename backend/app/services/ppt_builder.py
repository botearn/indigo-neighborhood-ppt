import io
import base64
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from app.core.models import StoryUnit, Beat, VisualIntent


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

C_BG = RGBColor(0x0F, 0x11, 0x16)
C_PANEL = RGBColor(0x13, 0x16, 0x1C)
C_DARK = RGBColor(0x22, 0x1A, 0x12)
C_CREAM = RGBColor(0xF5, 0xF0, 0xE6)
C_AMBER = RGBColor(0xC8, 0xA9, 0x6E)
C_AMBER_DIM = RGBColor(0x42, 0x30, 0x1C)
C_BODY = RGBColor(0xE8, 0xE2, 0xD4)
C_MUTED = RGBColor(0xA8, 0xA3, 0x97)
C_DIM = RGBColor(0x6B, 0x72, 0x80)
C_OVERLAY = RGBColor(0x08, 0x06, 0x04)


def _decode_data_url(data_url: str) -> bytes:
    _, b64 = data_url.split(",", 1)
    return base64.b64decode(b64)


def _fetch_image(url: str) -> bytes | None:
    if url.startswith("data:"):
        try:
            return _decode_data_url(url)
        except Exception:
            return None
    try:
        r = httpx.get(url, timeout=15, follow_redirects=True)
        r.raise_for_status()
        return r.content
    except Exception:
        return None


def build_ppt_from_slides(slide_data_urls: list[str]) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for url in slide_data_urls:
        slide = prs.slides.add_slide(blank)
        png = _decode_data_url(url)
        slide.shapes.add_picture(io.BytesIO(png), 0, 0, width=SLIDE_W, height=SLIDE_H)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _alpha(shape, val: int):
    try:
        sp = shape.fill._xPr
        sf = sp.find(qn("a:solidFill"))
        clr = sf.find(qn("a:srgbClr"))
        if clr is None:
            clr = sf.find(qn("a:sysClr"))
        a = etree.SubElement(clr, qn("a:alpha"))
        a.set("val", str(val))
    except Exception:
        pass


def _box(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    if alpha:
        _alpha(s, alpha)
    return s


def _txt(slide, text, l, t, w, h, *, sz=22, bold=False, color=C_CREAM, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = color


def _rule(slide, l, t, w, color=C_AMBER, thick=1.5):
    s = slide.shapes.add_shape(1, l, t, w, Pt(thick))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _picture(slide, url: str, l, t, w, h):
    img = _fetch_image(url)
    if not img:
        return False
    try:
        slide.shapes.add_picture(io.BytesIO(img), l, t, width=w, height=h)
        return True
    except Exception:
        return False


def _vbar(slide, l, t, h, color=C_AMBER, w=Pt(2.5)):
    s = slide.shapes.add_shape(1, l, t, int(w), h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _verb_eyebrow(slide, beat: Beat, index: int, total: int, l, t):
    _txt(slide, f"{index:02d} / {total:02d}", l, t, Inches(2), Inches(0.32),
         sz=10, color=C_DIM)
    _txt(slide, beat.verb, l, t + Inches(0.42), Inches(2), Inches(0.4),
         sz=12, bold=True, color=C_AMBER)


def _sensory_line(beat: Beat) -> str:
    return "  ·  ".join(d.description for d in beat.sensory)


def _cover_slide(prs, story: StoryUnit):
    s = _blank(prs)
    _bg(s)
    if story.mood_image_url and _picture(s, story.mood_image_url, 0, 0, SLIDE_W, SLIDE_H):
        _box(s, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x08, 0x06, 0x04), alpha=70000)
    _txt(s, "HOTEL INDIGO", Inches(1), Inches(0.7), Inches(11.3), Inches(0.4),
         sz=10, color=C_AMBER)
    _txt(s, story.signature.zh, Inches(1), Inches(3.0), Inches(11.3), Inches(2.0),
         sz=72, color=C_CREAM)
    _txt(s, story.signature.en.upper(), Inches(1), Inches(5.1), Inches(11.3), Inches(0.5),
         sz=18, color=C_AMBER)
    _txt(s, f"{story.city.upper()}  ·  {story.neighborhood.upper()}",
         Inches(1), Inches(6.7), Inches(11.3), Inches(0.4), sz=11, color=C_MUTED)


def _hook_slide(prs, story: StoryUnit):
    s = _blank(prs)
    _bg(s, C_PANEL)
    _txt(s, "STREET ENTRANCE", Inches(1.2), Inches(1.2), Inches(8), Inches(0.4),
         sz=10, color=C_AMBER)
    _rule(s, Inches(1.2), Inches(2.0), Inches(0.7))
    _txt(s, story.hook_line, Inches(1.2), Inches(2.6), Inches(11), Inches(2.5),
         sz=48, color=C_CREAM)
    _txt(s, story.anchor, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
         sz=18, color=C_MUTED)


def _layout_image_dominant(s, beat: Beat, index: int, total: int):
    """Full-bleed photo, gradient overlay pulling to bottom, text anchored low."""
    if not (beat.image_url and _picture(s, beat.image_url, 0, 0, SLIDE_W, SLIDE_H)):
        _box(s, 0, 0, SLIDE_W, SLIDE_H, C_PANEL)
    _box(s, 0, Inches(2.8), SLIDE_W, Inches(1.4), C_OVERLAY, alpha=35000)
    _box(s, 0, Inches(3.8), SLIDE_W, Inches(1.4), C_OVERLAY, alpha=55000)
    _box(s, 0, Inches(4.8), SLIDE_W, Inches(2.7), C_OVERLAY, alpha=82000)
    M = Inches(1.1)
    _verb_eyebrow(s, beat, index, total, M, Inches(3.55))
    _txt(s, beat.title, M, Inches(4.55), Inches(11), Inches(1.1),
         sz=36, color=C_CREAM)
    _rule(s, M, Inches(5.7), Inches(1.2))
    _txt(s, beat.copy, M, Inches(5.95), Inches(11), Inches(0.6),
         sz=18, color=C_BODY)
    if beat.detail:
        _txt(s, beat.detail, M, Inches(6.4), Inches(11), Inches(0.7),
             sz=12, color=C_MUTED)
    sensory = _sensory_line(beat)
    if sensory:
        _txt(s, sensory, M, Inches(7.05), Inches(11.5), Inches(0.4),
             sz=9, color=C_DIM)


def _layout_dense_detail(s, beat: Beat, index: int, total: int):
    """Text left, framed inset image right (evidence window)."""
    _bg(s)
    M = Inches(1.1)
    _verb_eyebrow(s, beat, index, total, M, Inches(1.35))
    _txt(s, beat.title, M, Inches(2.35), Inches(5), Inches(1.0),
         sz=34, color=C_CREAM)
    _rule(s, M, Inches(3.4), Inches(1.2))
    _txt(s, beat.copy, M, Inches(3.65), Inches(5), Inches(1.2),
         sz=17, color=C_BODY)
    if beat.detail:
        _txt(s, beat.detail, M, Inches(5.0), Inches(5), Inches(1.6),
             sz=12, color=C_MUTED)
    sensory = _sensory_line(beat)
    if sensory:
        _txt(s, sensory, M, Inches(6.85), Inches(5.2), Inches(0.4),
             sz=9, color=C_DIM)
    IL, IT, IW, IH = Inches(6.8), Inches(0.6), Inches(6.0), Inches(6.3)
    if not (beat.image_url and _picture(s, beat.image_url, IL, IT, IW, IH)):
        _box(s, IL, IT, IW, IH, C_PANEL)
    _rule(s, IL, IT, IW, thick=2.5)
    _rule(s, IL, IT + IH - Pt(2.5), IW, thick=2.5)
    _vbar(s, IL, IT, IH, w=Pt(2.5))
    _vbar(s, IL + IW - Pt(2.5), IT, IH, w=Pt(2.5))


def _layout_atmospheric(s, beat: Beat, index: int, total: int):
    """Full bleed + giant dim verb watermark left + content right half."""
    if not (beat.image_url and _picture(s, beat.image_url, 0, 0, SLIDE_W, SLIDE_H)):
        _box(s, 0, 0, SLIDE_W, SLIDE_H, C_PANEL)
    _box(s, 0, 0, SLIDE_W, SLIDE_H, C_OVERLAY, alpha=60000)
    _txt(s, beat.verb, Inches(0.5), Inches(0.9), Inches(7.0), Inches(4.0),
         sz=130, bold=True, color=C_AMBER_DIM)
    RL, RW = Inches(7.4), Inches(5.3)
    _verb_eyebrow(s, beat, index, total, RL, Inches(1.35))
    _txt(s, beat.title, RL, Inches(2.35), RW, Inches(1.0),
         sz=32, color=C_CREAM)
    _rule(s, RL, Inches(3.4), Inches(1.4))
    _txt(s, beat.copy, RL, Inches(3.65), RW, Inches(1.0),
         sz=16, color=C_BODY)
    if beat.detail:
        _txt(s, beat.detail, RL, Inches(4.85), RW, Inches(1.6),
             sz=12, color=C_MUTED)
    sensory = _sensory_line(beat)
    if sensory:
        _txt(s, sensory, RL, Inches(6.7), RW, Inches(0.4),
             sz=9, color=C_DIM)


def _layout_editorial_break(s, beat: Beat, index: int, total: int):
    """Horizontal split: image top, dark panel bottom, thick amber rule between."""
    HSPLIT = Inches(4.05)
    _bg(s)
    if not (beat.image_url and _picture(s, beat.image_url, 0, 0, SLIDE_W, HSPLIT)):
        _box(s, 0, 0, SLIDE_W, HSPLIT, C_PANEL)
    _box(s, 0, HSPLIT, SLIDE_W, SLIDE_H - HSPLIT, C_DARK)
    _rule(s, 0, HSPLIT - Pt(1.5), SLIDE_W, thick=3.0)
    M = Inches(1.1)
    _verb_eyebrow(s, beat, index, total, M, HSPLIT + Inches(0.35))
    _txt(s, beat.title, M, HSPLIT + Inches(1.35), Inches(11), Inches(0.95),
         sz=32, color=C_CREAM)
    _txt(s, beat.copy, M, HSPLIT + Inches(2.3), Inches(11), Inches(0.7),
         sz=16, color=C_BODY)
    if beat.detail:
        _txt(s, beat.detail, M, HSPLIT + Inches(3.05), Inches(11), Inches(1.0),
             sz=12, color=C_MUTED)


def _layout_typography_first(s, beat: Beat, index: int, total: int):
    """Words carry it. Image absent or as small thumbnail; centered punchy type."""
    _bg(s, C_DARK)
    M = Inches(1.1)
    _verb_eyebrow(s, beat, index, total, M, Inches(0.9))
    # Optional small evidence thumbnail upper-right
    if beat.image_url:
        _picture(s, beat.image_url, Inches(10.3), Inches(0.7), Inches(2.4), Inches(1.6))
    # Massive title centered
    _txt(s, beat.title, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.6),
         sz=60, color=C_CREAM, align=PP_ALIGN.CENTER)
    _rule(s, Inches(6.17), Inches(4.15), Inches(1.0))
    _txt(s, beat.copy, Inches(0.8), Inches(4.45), Inches(11.7), Inches(1.0),
         sz=22, color=C_AMBER, align=PP_ALIGN.CENTER)
    if beat.detail:
        _txt(s, beat.detail, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.4),
             sz=13, color=C_MUTED, align=PP_ALIGN.CENTER)
    sensory = _sensory_line(beat)
    if sensory:
        _txt(s, sensory, Inches(1), Inches(7.0), Inches(11.3), Inches(0.4),
             sz=9, color=C_DIM, align=PP_ALIGN.CENTER)


def _layout_quiet_balance(s, beat: Beat, index: int, total: int):
    """Vertical split: dark text panel left, raw image right with amber divider."""
    _bg(s)
    SPLIT = Inches(5.6)
    if not (beat.image_url and _picture(s, beat.image_url, SPLIT, 0, SLIDE_W - SPLIT, SLIDE_H)):
        _box(s, SPLIT, 0, SLIDE_W - SPLIT, SLIDE_H, C_PANEL)
    _vbar(s, SPLIT - Pt(1), 0, SLIDE_H, color=C_AMBER, w=Pt(2))
    M = Inches(1.1)
    _txt(s, f"{index:02d}", M, Inches(0.65), Inches(2), Inches(0.9),
         sz=40, bold=True, color=C_AMBER_DIM)
    _txt(s, beat.verb, M, Inches(1.7), Inches(4), Inches(0.4),
         sz=12, bold=True, color=C_AMBER)
    _txt(s, beat.title, M, Inches(2.25), Inches(4.3), Inches(1.1),
         sz=32, color=C_CREAM)
    _rule(s, M, Inches(3.5), Inches(1.2))
    _txt(s, beat.copy, M, Inches(3.75), Inches(4.3), Inches(1.0),
         sz=16, color=C_BODY)
    if beat.detail:
        _txt(s, beat.detail, M, Inches(4.95), Inches(4.3), Inches(1.6),
             sz=12, color=C_MUTED)
    sensory = _sensory_line(beat)
    if sensory:
        _txt(s, sensory, M, Inches(6.85), Inches(4.4), Inches(0.4),
             sz=9, color=C_DIM)


_LAYOUTS = {
    VisualIntent.IMAGE_DOMINANT: _layout_image_dominant,
    VisualIntent.TYPOGRAPHY_FIRST: _layout_typography_first,
    VisualIntent.QUIET_BALANCE: _layout_quiet_balance,
    VisualIntent.DENSE_DETAIL: _layout_dense_detail,
    VisualIntent.ATMOSPHERIC: _layout_atmospheric,
    VisualIntent.EDITORIAL_BREAK: _layout_editorial_break,
}


def _infer_intent(beat: Beat) -> VisualIntent:
    """Fallback when LLM omits visual_intent — derive from content signals."""
    has_image = bool(beat.image_url)
    detail_len = len(beat.detail or "")
    sensory_count = len(beat.sensory)
    if not has_image:
        return VisualIntent.TYPOGRAPHY_FIRST
    if detail_len >= 80:
        return VisualIntent.DENSE_DETAIL
    if sensory_count >= 4:
        return VisualIntent.ATMOSPHERIC
    return VisualIntent.IMAGE_DOMINANT


def _beat_slide(prs, beat: Beat, index: int, total: int):
    s = _blank(prs)
    _bg(s)
    intent = beat.visual_intent or _infer_intent(beat)
    layout = _LAYOUTS.get(intent, _layout_quiet_balance)
    layout(s, beat, index, total)


def _action_slide(prs, story: StoryUnit):
    s = _blank(prs)
    _bg(s)
    _txt(s, "STEP OUT", Inches(1), Inches(2.0), Inches(11.3), Inches(0.4),
         sz=11, color=C_AMBER, align=PP_ALIGN.CENTER)
    _rule(s, Inches(6.17), Inches(2.7), Inches(1.0))
    _txt(s, story.action_cue, Inches(1), Inches(3.2), Inches(11.3), Inches(2),
         sz=44, color=C_CREAM, align=PP_ALIGN.CENTER)
    _rule(s, Inches(6.17), Inches(5.3), Inches(1.0))


def _closing_slide(prs, story: StoryUnit):
    s = _blank(prs)
    _bg(s)
    _txt(s, story.signature.zh, Inches(1), Inches(2.7), Inches(11.3), Inches(1.5),
         sz=56, color=C_CREAM, align=PP_ALIGN.CENTER)
    _txt(s, story.signature.en.upper(), Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
         sz=14, color=C_AMBER, align=PP_ALIGN.CENTER)
    _txt(s, f"HOTEL INDIGO  ·  {story.city.upper()}  ·  {story.neighborhood.upper()}",
         Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
         sz=10, color=C_DIM, align=PP_ALIGN.CENTER)


def build_ppt_from_story(story: StoryUnit) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _cover_slide(prs, story)
    _hook_slide(prs, story)
    total = len(story.beats)
    for i, beat in enumerate(story.beats, start=1):
        _beat_slide(prs, beat, i, total)
    _action_slide(prs, story)
    _closing_slide(prs, story)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
