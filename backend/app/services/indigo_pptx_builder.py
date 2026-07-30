"""
Builds a fully-editable 22-slide PPTX from IndigoStoryUnit using python-pptx.
All text elements are real text boxes (not screenshots), fonts and colors are
applied programmatically. Image areas are placeholder shapes.
"""
import base64
from concurrent.futures import ThreadPoolExecutor, as_completed
from contextvars import ContextVar
from io import BytesIO
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from lxml import etree

from app.core.models import IndigoStoryUnit, IndigoBeat

_IMAGE_RAW_CACHE: ContextVar[dict[str, bytes | None] | None] = ContextVar(
    "indigo_pptx_image_raw_cache",
    default=None,
)
_IMAGE_JPEG_CACHE: ContextVar[dict[tuple[str, float], bytes | None] | None] = ContextVar(
    "indigo_pptx_image_jpeg_cache",
    default=None,
)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────
SW = Cm(33.867)
SH = Cm(19.05)

# ── Color tokens ─────────────────────────────────────────────────────────
TEAL   = RGBColor(0x2D, 0x7A, 0x7A)
TEAL_L = RGBColor(0x3A, 0x9A, 0x9A)
GOLD   = RGBColor(0xC8, 0xA9, 0x6E)
NAVY   = RGBColor(0x1A, 0x2E, 0x3B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_D = RGBColor(0x37, 0x41, 0x51)
GRAY_M = RGBColor(0x6B, 0x72, 0x80)
GRAY_L = RGBColor(0xD1, 0xD5, 0xDB)

# Dark slide backgrounds per type
BG_DARK   = RGBColor(0x0F, 0x18, 0x20)   # cover / cinematic
BG_CINEMA = RGBColor(0x0E, 0x16, 0x10)   # story emotion
BG_SUMM   = RGBColor(0x12, 0x18, 0x0A)   # story summary

BEAT_BG = {
    "01": RGBColor(0x1E, 0x14, 0x08),
    "02": RGBColor(0x0E, 0x18, 0x18),
    "03": RGBColor(0x1E, 0x16, 0x0A),
    "04": RGBColor(0x1A, 0x0E, 0x0A),
    "05": RGBColor(0x0E, 0x12, 0x18),
    "06": RGBColor(0x0E, 0x1A, 0x18),
}
ORIGIN_BG = [
    RGBColor(0x2A, 0x3E, 0x3E),
    RGBColor(0x3A, 0x30, 0x20),
    RGBColor(0x28, 0x3A, 0x28),
]


# ── Helpers ───────────────────────────────────────────────────────────────

def _new_slide(prs: Presentation) -> object:
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
              line_color: RGBColor | None = None, line_width: float = 0):
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()  # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def _add_oval(slide, left, top, width, height, fill: RGBColor,
              line_color: RGBColor | None = None, line_width: float = 0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def _add_text(slide, text: str, left, top, width, height,
              size: float = 10,
              bold: bool = False,
              color: RGBColor = WHITE,
              align: PP_ALIGN = PP_ALIGN.LEFT,
              wrap: bool = True,
              font: str = "Helvetica Neue",
              spacing: float | None = None,
              line_spacing: float | None = None,
              vertical_anchor: str = "top") -> object:
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap

    anchors = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as PtU
        p.line_spacing = Pt(line_spacing)

    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    if spacing:
        # letter spacing via XML
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return txBox


def _add_multiline(slide, lines: list[str], left, top, width, height,
                   size: float = 10, bold: bool = False, color: RGBColor = WHITE,
                   align: PP_ALIGN = PP_ALIGN.LEFT, font: str = "Helvetica Neue",
                   line_gap: float = 1.4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = Pt(size * line_gap)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return txBox


def _hbar(slide, city: str, district: str):
    """4-cell header bar, 1.2cm tall, teal bottom border."""
    bar_h = Cm(1.2)
    _add_rect(slide, 0, 0, SW, bar_h, fill=RGBColor(0xFF, 0xFF, 0xFF))
    _add_rect(slide, 0, bar_h - Cm(0.07), SW, Cm(0.07), fill=TEAL)
    col_w = SW // 4
    labels = [
        "HOTELINDIGO",
        f"HOTEL INDIGO {city.upper()} {district.upper()}  |  STORYLINE STRATEGIC DEVELOPMENT",
        f"英迪格酒店 · {city}{district}  |  2026",
    ]
    for i, label in enumerate(labels):
        _add_text(slide, label,
                  left=col_w * i + Cm(0.3), top=Cm(0.05),
                  width=col_w - Cm(0.3), height=bar_h,
                  size=5.5, color=NAVY, font="Helvetica Neue",
                  vertical_anchor="middle")


def _sec_label(slide, en: str, zh: str):
    _add_text(slide, en, Cm(0.8), Cm(1.35), Cm(8), Cm(0.5),
              size=5.5, bold=True, color=TEAL, spacing=0.2)
    _add_text(slide, zh, Cm(0.8), Cm(1.75), Cm(8), Cm(0.5),
              size=5, color=GRAY_M)


def _page_num(slide, n: int, dark: bool = False):
    color = GRAY_M if dark else RGBColor(0xA0, 0xA0, 0x98)
    _add_text(slide, str(n),
              left=SW - Cm(1.5), top=SH - Cm(0.8),
              width=Cm(1.2), height=Cm(0.6),
              size=9, color=color, align=PP_ALIGN.RIGHT)


def _img_placeholder(slide, left, top, width, height, label: str = "",
                     bg: RGBColor = RGBColor(0xCC, 0xD4, 0xD4),
                     image_url: str | None = None):
    if image_url:
        img_stream = _resolve_image(image_url, width, height)
        if img_stream:
            slide.shapes.add_picture(img_stream, left, top, width, height)
            return
    _add_rect(slide, left, top, width, height, fill=bg)
    if label:
        _add_text(slide, f"[ {label} ]",
                  left + Cm(0.3), top + height - Cm(0.7),
                  width - Cm(0.6), Cm(0.6),
                  size=5, color=RGBColor(0xAA, 0xAA, 0xAA))


def _beat_image(s: IndigoStoryUnit, idx: int, *fields: str) -> str | None:
    if not s.beats:
        return None
    beat = s.beats[idx % len(s.beats)]
    image_fields = fields or ("image_url", "mood_image_url", "col2_image_url", "col3_image_url")
    for field in image_fields:
        url = getattr(beat, field, None)
        if url:
            return url
    return None


def _center_crop_to_ratio(img, target_ratio: float):
    width, height = img.size
    if width <= 0 or height <= 0 or target_ratio <= 0:
        return img

    current_ratio = width / height
    if abs(current_ratio - target_ratio) < 0.01:
        return img

    if current_ratio > target_ratio:
        new_width = int(height * target_ratio)
        left = max(0, (width - new_width) // 2)
        return img.crop((left, 0, left + new_width, height))

    new_height = int(width / target_ratio)
    top = max(0, (height - new_height) // 2)
    return img.crop((0, top, width, top + new_height))


def _load_raw_image(url: str) -> bytes | None:
    if url.startswith("data:"):
        _, payload = url.split(",", 1)
        return base64.b64decode(payload)
    if url.startswith("http"):
        import httpx
        response = httpx.get(url, timeout=10, follow_redirects=True)
        response.raise_for_status()
        return response.content
    return None


def _story_image_urls(story: IndigoStoryUnit) -> list[str]:
    urls = (
        getattr(beat, field, None)
        for beat in story.beats
        for field in ("image_url", "mood_image_url", "col2_image_url", "col3_image_url")
    )
    return list(dict.fromkeys(url for url in urls if url))


def _prefetch_story_images(story: IndigoStoryUnit, cache: dict[str, bytes | None]) -> None:
    urls = _story_image_urls(story)
    if not urls:
        return

    worker_count = min(8, len(urls))
    with ThreadPoolExecutor(max_workers=worker_count, thread_name_prefix="pptx-image") as executor:
        future_urls = {executor.submit(_load_raw_image, url): url for url in urls}
        for future in as_completed(future_urls):
            url = future_urls[future]
            try:
                cache[url] = future.result()
            except Exception:
                cache[url] = None


def _resolve_image(url: str, frame_width: Emu | None = None, frame_height: Emu | None = None) -> BytesIO | None:
    """Turn a data:…;base64 URL or http(s) URL into a JPEG BytesIO for python-pptx."""
    if not url:
        return None
    raw_cache = _IMAGE_RAW_CACHE.get()
    jpeg_cache = _IMAGE_JPEG_CACHE.get()
    target_ratio = round(float(frame_width) / float(frame_height), 4) if frame_width and frame_height else 0.0
    jpeg_cache_key = (url, target_ratio)
    try:
        if jpeg_cache is not None and jpeg_cache_key in jpeg_cache:
            cached_jpeg = jpeg_cache[jpeg_cache_key]
            return BytesIO(cached_jpeg) if cached_jpeg else None

        if raw_cache is not None and url in raw_cache:
            raw = raw_cache[url]
        else:
            raw = _load_raw_image(url)
            if raw_cache is not None:
                raw_cache[url] = raw
        if raw:
            from PIL import Image, ImageOps
            img = Image.open(BytesIO(raw))
            img = ImageOps.exif_transpose(img)
            if frame_width and frame_height:
                img = _center_crop_to_ratio(img, target_ratio)
            buf = BytesIO()
            img.convert("RGB").save(buf, format="JPEG", quality=85)
            jpeg = buf.getvalue()
            if jpeg_cache is not None:
                jpeg_cache[jpeg_cache_key] = jpeg
            return BytesIO(jpeg)
    except Exception:
        if jpeg_cache is not None:
            jpeg_cache[jpeg_cache_key] = None
    return None


def _split_origin_body(text: str) -> tuple[str, str]:
    if len(text) < 40:
        return text, ""
    midpoint = len(text) // 2
    candidates = [i for i, ch in enumerate(text) if ch in "。；，、" and 12 < i < len(text) - 12]
    if not candidates:
        return text[:midpoint], text[midpoint:]
    split_at = min(candidates, key=lambda i: abs(i - midpoint)) + 1
    return text[:split_at], text[split_at:]


# ── Slide builders ────────────────────────────────────────────────────────

def _slide01_cover(prs, s: IndigoStoryUnit):
    slide = _new_slide(prs)
    _set_bg(slide, BG_DARK)
    _hbar(slide, s.city, s.district)

    cover_image = _beat_image(s, 0, "image_url", "mood_image_url", "col2_image_url", "col3_image_url")
    if cover_image:
        _img_placeholder(slide, SW * 0.52, Cm(1.2), SW * 0.48, SH - Cm(1.2),
                         bg=BG_DARK, image_url=cover_image)
        _add_rect(slide, SW * 0.50, Cm(1.2), Cm(0.25), SH - Cm(1.2), fill=BG_DARK)

    _add_text(slide, "HOTEL", Cm(1.5), SH - Cm(6.5), Cm(15), Cm(1),
              size=8, color=RGBColor(0x80, 0x80, 0x80), spacing=0.3)
    _add_text(slide, "INDIGO", Cm(1.5), SH - Cm(5.8), Cm(20), Cm(2.2),
              size=44, bold=True, color=WHITE, spacing=0.12, font="Helvetica Neue")
    _add_text(slide, f"Hotel Indigo {s.hotel_en}  ·  Touchpoints Development",
              Cm(1.5), SH - Cm(3.4), Cm(22), Cm(0.7),
              size=10, color=RGBColor(0x99, 0x99, 0x99))
    _add_text(slide, "PHASE 01  ·  PROPOSAL  ·  2026",
              Cm(1.5), SH - Cm(2.7), Cm(15), Cm(0.6),
              size=7.5, color=TEAL_L, spacing=0.15)
    _page_num(slide, 1)


def _slide02_taglines(prs, s: IndigoStoryUnit):
    slide = _new_slide(prs)
    _set_bg(slide, WHITE)
    _hbar(slide, s.city, s.district)
    _sec_label(slide, "TAGLINE OPTION", "故事标题方案")

    col_w = (SW - Cm(2.0)) / 3
    for i, tl in enumerate(s.taglines):
        x = Cm(0.7) + col_w * i + Cm(0.2) * i
        border = TEAL if i == 0 else GRAY_L
        image_url = _beat_image(s, i, "mood_image_url", "image_url", "col2_image_url", "col3_image_url")
        _add_rect(slide, x, Cm(2.3), col_w - Cm(0.15), SH - Cm(2.8),
                  fill=RGBColor(0xFA, 0xFF, 0xFF) if i == 0 else WHITE,
                  line_color=border, line_width=0.75)
        _add_text(slide, f"Option {i+1}{'  ★' if i == 0 else ''}",
                  x + Cm(0.4), Cm(2.7), col_w - Cm(0.8), Cm(0.5),
                  size=5.5, color=GRAY_M, spacing=0.12)
        _add_text(slide, tl.zh,
                  x + Cm(0.4), Cm(3.3), col_w - Cm(0.8), Cm(2.0),
                  size=20, color=NAVY, spacing=0.5,
                  font="Songti SC")
        _add_text(slide, tl.sub,
                  x + Cm(0.4), Cm(5.5), col_w - Cm(0.8), Cm(3.0),
                  size=8.5, color=GRAY_M, wrap=True, line_spacing=14)
        _img_placeholder(slide, x + Cm(0.4), Cm(10.0), col_w - Cm(0.95), Cm(5.7),
                         bg=RGBColor(0xEC, 0xF0, 0xF0), image_url=image_url)
    _page_num(slide, 2, dark=True)


def _slide_cinematic(prs, n: int, s: IndigoStoryUnit, bg: RGBColor,
                     headline: str, paras: list[str], top_label: str,
                     image_url: str | None = None):
    slide = _new_slide(prs)
    _set_bg(slide, bg)
    if image_url:
        image_left = SW - Cm(9.6)
        _img_placeholder(slide, image_left, Cm(1.2), Cm(9.6), SH - Cm(1.2),
                         bg=bg, image_url=image_url)
        _add_rect(slide, image_left - Cm(0.25), Cm(1.2), Cm(0.25), SH - Cm(1.2), fill=bg)
    _add_text(slide, top_label, Cm(0.8), Cm(0.6), Cm(10), Cm(0.5),
              size=6, color=RGBColor(0x80, 0x80, 0x80), spacing=0.2)
    _add_text(slide, "故事概念方向", Cm(0.8), Cm(1.0), Cm(10), Cm(0.5),
              size=5.5, color=RGBColor(0x55, 0x55, 0x55))
    text_left = Cm(2.6)
    text_width = SW - Cm(15) if image_url else SW - Cm(6)
    # Headline centered
    _add_text(slide, headline,
              text_left, SH * 0.28, text_width, Cm(3),
              size=17, color=WHITE, align=PP_ALIGN.CENTER,
              font="Songti SC", wrap=True, line_spacing=28)
    # Paragraphs
    for i, para in enumerate(paras):
        _add_text(slide, para,
                  text_left + Cm(1), SH * 0.52 + Cm(2.2) * i, text_width - Cm(2), Cm(2),
                  size=8, color=RGBColor(0xAA, 0xAA, 0xAA),
                  align=PP_ALIGN.CENTER, wrap=True, line_spacing=14)
    _page_num(slide, n)


def _slide_origin(prs, n: int, s: IndigoStoryUnit, idx: int):
    slide = _new_slide(prs)
    o = s.origins[idx]
    beat_img = _beat_image(s, idx, "image_url", "mood_image_url", "col2_image_url", "col3_image_url")

    if idx == 0:
        _set_bg(slide, WHITE)
        _hbar(slide, s.city, s.district)
        _sec_label(slide, "STORYLINE CONCEPT", "故事概念方向")
        photo_w = Cm(12.4)
        _img_placeholder(slide, 0, Cm(1.2), photo_w, SH - Cm(1.2),
                         label=o.title, bg=ORIGIN_BG[idx], image_url=beat_img)
        _add_rect(slide, photo_w, Cm(1.2), Cm(0.08), SH - Cm(1.2), fill=TEAL)
        tx = photo_w + Cm(1.0)
        tw = SW - tx - Cm(1.2)
        _add_text(slide, f"ORIGIN 0{idx + 1}  ·  {o.title}",
                  tx, Cm(2.15), tw, Cm(0.55),
                  size=6, color=TEAL, spacing=0.18)
        _add_text(slide, o.headline,
                  tx, Cm(3.1), tw - Cm(1.2), Cm(2.1),
                  size=14, bold=True, color=NAVY, wrap=True, line_spacing=22,
                  font="Songti SC")
        _add_text(slide, o.body,
                  tx, Cm(5.9), tw, Cm(8.6),
                  size=8, color=GRAY_D, wrap=True, line_spacing=14)
        _add_text(slide, s.story_summary,
                  tx, SH - Cm(2.5), tw - Cm(2.2), Cm(0.8),
                  size=6.5, color=TEAL, wrap=True)
        _page_num(slide, n, dark=True)
        return

    if idx == 1:
        _set_bg(slide, WHITE)
        _hbar(slide, s.city, s.district)
        _img_placeholder(slide, 0, Cm(1.2), SW, Cm(7.0),
                         label=o.title, bg=ORIGIN_BG[idx], image_url=beat_img)
        _add_rect(slide, Cm(1.2), Cm(6.85), SW - Cm(2.4), Cm(3.1), fill=WHITE)
        _add_text(slide, f"ORIGIN 0{idx + 1}  ·  {o.title}",
                  Cm(2.0), Cm(7.25), Cm(9), Cm(0.55),
                  size=6, color=TEAL, spacing=0.18)
        _add_text(slide, o.headline,
                  Cm(2.0), Cm(8.05), SW - Cm(4), Cm(1.5),
                  size=14.5, bold=True, color=NAVY, wrap=True, line_spacing=22,
                  font="Songti SC")
        left_body, right_body = _split_origin_body(o.body)
        _add_text(slide, left_body,
                  Cm(2.0), Cm(11.1), Cm(13.6), Cm(4.2),
                  size=7.8, color=GRAY_D, wrap=True, line_spacing=13)
        _add_text(slide, right_body,
                  Cm(17.2), Cm(11.1), Cm(13.8), Cm(4.2),
                  size=7.8, color=GRAY_D, wrap=True, line_spacing=13)
        _add_text(slide, "NEIGHBORHOOD MEMORY",
                  Cm(2.0), SH - Cm(2.0), Cm(12), Cm(0.55),
                  size=6, color=GRAY_M, spacing=0.22)
        _page_num(slide, n, dark=True)
        return

    _set_bg(slide, ORIGIN_BG[idx])
    _hbar(slide, s.city, s.district)
    image_w = Cm(13.8)
    image_x = SW - image_w
    _img_placeholder(slide, image_x, Cm(1.2), image_w, SH - Cm(1.2),
                     label=o.title, bg=ORIGIN_BG[idx], image_url=beat_img)
    _add_rect(slide, image_x - Cm(0.18), Cm(1.2), Cm(0.18), SH - Cm(1.2), fill=ORIGIN_BG[idx])
    _add_text(slide, f"ORIGIN 0{idx + 1}",
              Cm(1.2), Cm(2.1), Cm(7), Cm(0.5),
              size=6, color=TEAL_L, spacing=0.22)
    _add_text(slide, o.title,
              Cm(1.2), Cm(2.75), Cm(9), Cm(0.8),
              size=8, color=GOLD, spacing=0.12)
    _add_text(slide, o.headline,
              Cm(1.2), Cm(4.35), image_x - Cm(2.6), Cm(3.0),
              size=18, bold=True, color=WHITE, wrap=True, line_spacing=27,
              font="Songti SC")
    _add_text(slide, o.body,
              Cm(1.2), Cm(8.3), image_x - Cm(3.0), Cm(5.8),
              size=8, color=RGBColor(0xD6, 0xD6, 0xCC), wrap=True, line_spacing=14)
    _add_text(slide, f"{s.city}  ·  {s.district}",
              Cm(1.2), SH - Cm(2.1), Cm(10), Cm(0.6),
              size=6.5, color=TEAL_L, spacing=0.18)
    _page_num(slide, n)


def _slide_story_summary(prs, n: int, s: IndigoStoryUnit):
    slide = _new_slide(prs)
    _set_bg(slide, BG_SUMM)
    image_url = _beat_image(s, 4, "mood_image_url", "image_url", "col2_image_url", "col3_image_url")
    text_width = SW
    if image_url:
        image_w = Cm(10.2)
        _img_placeholder(slide, SW - image_w, Cm(1.1), image_w, SH - Cm(1.1),
                         bg=BG_SUMM, image_url=image_url)
        _add_rect(slide, SW - image_w - Cm(0.25), Cm(1.1), Cm(0.25), SH - Cm(1.1), fill=BG_SUMM)
        text_width = SW - image_w - Cm(0.5)
    _add_text(slide, "STORY SUMMARY", Cm(0.8), Cm(0.6), Cm(10), Cm(0.5),
              size=6, color=GRAY_M, spacing=0.2)
    _add_text(slide, f"Hotel Indigo {s.hotel_en}",
              Cm(0), SH * 0.32, text_width, Cm(0.8),
              size=7, color=TEAL_L, spacing=0.25, align=PP_ALIGN.CENTER)
    _add_text(slide, s.story_summary,
              Cm(2), SH * 0.40, text_width - Cm(4), Cm(3.5),
              size=13, color=WHITE, align=PP_ALIGN.CENTER, wrap=True,
              line_spacing=22, font="Songti SC")
    _add_text(slide, f"{s.city.upper()}  ·  {s.district.upper()}",
              Cm(0), SH * 0.75, text_width, Cm(0.6),
              size=7, color=GRAY_M, spacing=0.15, align=PP_ALIGN.CENTER)
    _page_num(slide, n)


def _slide_story_mapping(prs, n: int, s: IndigoStoryUnit):
    slide = _new_slide(prs)
    _set_bg(slide, WHITE)
    _add_text(slide, "STORY MAPPING", Cm(1), Cm(1.3), Cm(12), Cm(1.0),
              size=18, bold=True, color=NAVY, spacing=0.04)
    _add_text(slide, "空间触点索引", Cm(13.5), Cm(1.7), Cm(6), Cm(0.7),
              size=9, color=GRAY_M)
    col_w = (SW - Cm(2)) / 6
    colors_bg = [BEAT_BG[b.num] for b in s.beats]
    for i, (beat, bg) in enumerate(zip(s.beats, colors_bg)):
        x = Cm(1) + col_w * i
        th_h = SH - Cm(5.5)
        image_url = _beat_image(s, i, "image_url", "mood_image_url", "col2_image_url", "col3_image_url")
        _img_placeholder(slide, x, Cm(3.2), col_w - Cm(0.2), th_h,
                         bg=bg, image_url=image_url)
        _add_text(slide, beat.num, x + Cm(0.25), Cm(3.4), Cm(1.5), Cm(0.8),
                  size=11, bold=True, color=RGBColor(0xBB, 0xBB, 0xBB))
        _add_text(slide, beat.name_zh,
                  x, Cm(3.2) + th_h + Cm(0.2), col_w - Cm(0.1), Cm(0.9),
                  size=7, bold=True, color=NAVY, align=PP_ALIGN.CENTER, wrap=True)
        _add_text(slide, beat.space_zh,
                  x, Cm(3.2) + th_h + Cm(1.1), col_w - Cm(0.1), Cm(0.6),
                  size=5.5, color=GRAY_M, align=PP_ALIGN.CENTER, wrap=True)
    _page_num(slide, n, dark=True)


def _slide_story_flow_grid(prs, n: int, s: IndigoStoryUnit):
    slide = _new_slide(prs)
    _set_bg(slide, WHITE)
    _add_text(slide, "STORY TOUCHPOINTS", Cm(1), Cm(0.65), Cm(8), Cm(0.45),
              size=6, bold=True, color=TEAL, spacing=0.18)
    _add_text(slide, "空间触点", Cm(1), Cm(1.15), Cm(8), Cm(0.9),
              size=17, color=NAVY, spacing=0.06, font="Songti SC")
    _add_text(slide, s.story_summary,
              Cm(16.4), Cm(0.95), SW - Cm(17.4), Cm(1.15),
              size=7.5, color=GRAY_M, wrap=True, align=PP_ALIGN.RIGHT)

    margin_x = Cm(1)
    col_gap = Cm(0.65)
    row_gap = Cm(0.95)
    grid_top = Cm(3.15)
    grid_bottom = SH - Cm(0.9)
    col_w = (SW - margin_x * 2 - col_gap * 2) / 3
    row_h = (grid_bottom - grid_top - row_gap) / 2
    img_h = Cm(3.05)

    for i, beat in enumerate(s.beats):
        col = i % 3
        row = i // 3
        x = margin_x + (col_w + col_gap) * col
        y = grid_top + (row_h + row_gap) * row
        image_url = _beat_image(s, i, "image_url", "mood_image_url", "col2_image_url", "col3_image_url")
        _img_placeholder(slide, x, y, col_w, img_h,
                         bg=RGBColor(0xEC, 0xF0, 0xF0), image_url=image_url)
        _add_text(slide, beat.num,
                  x, y + Cm(3.35), Cm(0.9), Cm(0.45),
                  size=6.5, bold=True, color=TEAL, spacing=0.12)
        _add_text(slide, beat.name_zh,
                  x + Cm(1.0), y + Cm(3.25), col_w - Cm(1.0), Cm(0.65),
                  size=9, bold=True, color=NAVY, font="Songti SC", wrap=True)
        _add_text(slide, beat.space_zh,
                  x, y + Cm(4.02), col_w, Cm(0.45),
                  size=5.5, color=GRAY_M, spacing=0.08, wrap=True)
        _add_text(slide, beat.narrative[:58] + "…",
                  x, y + Cm(4.65), col_w, Cm(1.35),
                  size=6.2, color=GRAY_D, wrap=True, line_spacing=9.5)
        _add_rect(slide, x, y + row_h - Cm(0.72), Cm(0.75), Cm(0.035), fill=GOLD)
        _add_text(slide, beat.tagline,
                  x + Cm(0.95), y + row_h - Cm(0.86), col_w - Cm(0.95), Cm(0.45),
                  size=5.8, color=TEAL, wrap=True)
    _page_num(slide, n, dark=True)


def _slide_beat_cover(prs, n: int, beat: IndigoBeat):
    slide = _new_slide(prs)
    _set_bg(slide, BEAT_BG[beat.num])
    panel_w = SW * 0.44
    ghost_lines = beat.ghost_en.split("\n")

    # Beat title
    _add_text(slide, beat.name_zh,
              Cm(1.5), SH * 0.35, panel_w - Cm(2), Cm(2.0),
              size=24, color=WHITE, spacing=0.55,
              font="Songti SC")
    _add_text(slide, beat.space_zh,
              Cm(1.5), SH * 0.35 + Cm(2.2), panel_w - Cm(2), Cm(0.5),
              size=7.5, color=RGBColor(0xAA, 0xAA, 0xAA), spacing=0.07)
    # Ghost EN
    _add_multiline(slide, ghost_lines,
                   Cm(1.5), SH * 0.35 + Cm(3.0), panel_w - Cm(2), Cm(2.5),
                   size=24, bold=True, color=GOLD, font="Helvetica Neue", line_gap=1.1)
    # Narrative
    _add_text(slide, beat.narrative,
              Cm(1.5), SH * 0.35 + Cm(5.2), panel_w - Cm(2), Cm(2.5),
              size=8.5, color=RGBColor(0xBB, 0xBB, 0xBB), wrap=True, line_spacing=14)
    _add_text(slide, "在这里  ·  ZAI ZHE LI",
              Cm(1.5), SH - Cm(2.8), panel_w - Cm(2), Cm(0.5),
              size=7, color=TEAL_L, spacing=0.18)
    _add_text(slide, beat.tagline,
              Cm(1.5), SH - Cm(2.2), panel_w - Cm(2), Cm(0.6),
              size=10.5, bold=True, color=WHITE)

    # Image placeholder (right side)
    _img_placeholder(slide, panel_w + Cm(0.3), Cm(0.5),
                     SW - panel_w - Cm(0.3), SH - Cm(1),
                     label="Photo placeholder",
                     bg=RGBColor(0x22, 0x22, 0x22),
                     image_url=beat.image_url)
    _page_num(slide, n)


def _slide_moodboard(prs, n: int, beat: IndigoBeat, s: IndigoStoryUnit):
    slide = _new_slide(prs)
    _set_bg(slide, WHITE)
    _hbar(slide, s.city, s.district)
    _sec_label(slide, "STORYLINE CONCEPT", "故事概念方向")

    bar_y = Cm(2.05)
    c1_w = SW * 0.28
    c2_w = SW * 0.36
    c3_w = SW - c1_w - c2_w
    c2_x = c1_w
    c3_x = c1_w + c2_w

    # Col 1 separator lines
    _add_rect(slide, c1_w, bar_y, Cm(0.02), SH - bar_y, fill=GRAY_L)
    _add_rect(slide, c2_x + c2_w, bar_y, Cm(0.02), SH - bar_y, fill=GRAY_L)

    # Col 1: ghost + concept
    ghost_lines = beat.mb_ghost_en.split("\n")
    _add_multiline(slide, ghost_lines,
                   Cm(0.6), bar_y + Cm(0.4), c1_w - Cm(0.8), Cm(2.2),
                   size=17, bold=True, color=RGBColor(0xE0, 0xE8, 0xE8),
                   font="Helvetica Neue", line_gap=1.0)
    _add_text(slide, beat.space_zh,
              Cm(0.6), bar_y + Cm(2.8), c1_w - Cm(0.8), Cm(0.4),
              size=6, bold=True, color=NAVY, spacing=0.07)
    _add_text(slide, "运用元素",
              Cm(0.6), bar_y + Cm(3.2), c1_w - Cm(0.8), Cm(0.4),
              size=6, color=TEAL, spacing=0.15)
    _add_text(slide, beat.mb_concept,
              Cm(0.6), bar_y + Cm(3.7), c1_w - Cm(0.8), Cm(0.9),
              size=12, bold=True, color=NAVY, font="Songti SC")
    _add_text(slide, beat.mb_concept_sub,
              Cm(0.6), bar_y + Cm(4.6), c1_w - Cm(0.8), Cm(0.5),
              size=7, color=GRAY_M, wrap=True)
    # Col 1 image placeholder
    img_top = bar_y + Cm(5.2)
    _img_placeholder(slide, Cm(0.5), img_top, c1_w - Cm(0.8), SH - img_top - Cm(0.3),
                     bg=RGBColor(0xC8, 0xD4, 0xD4),
                     image_url=beat.mood_image_url)

    # Col 2
    _add_text(slide, beat.mb_col2_title,
              c2_x + Cm(0.5), bar_y + Cm(0.4), c2_w - Cm(0.8), Cm(0.7),
              size=8, bold=True, color=NAVY, wrap=True, spacing=0.04)
    _add_text(slide, beat.mb_col2_accent,
              c2_x + Cm(0.5), bar_y + Cm(1.1), c2_w - Cm(0.8), Cm(0.5),
              size=8, bold=True, color=TEAL, spacing=0.02)
    _add_text(slide, beat.mb_col2_body,
              c2_x + Cm(0.5), bar_y + Cm(1.7), c2_w - Cm(0.8), Cm(3.5),
              size=7, color=GRAY_D, wrap=True, line_spacing=12)
    # 2+1 image grid
    img_y = bar_y + Cm(5.4)
    img_h1 = Cm(2.4)
    img_h2 = Cm(1.7)
    _img_placeholder(slide, c2_x + Cm(0.5), img_y, (c2_w - Cm(1.2)) / 2, img_h1,
                     bg=RGBColor(0xCC, 0xCC, 0xBB), image_url=beat.col2_image_url)
    _img_placeholder(slide, c2_x + Cm(0.5) + (c2_w - Cm(1.2)) / 2 + Cm(0.15), img_y,
                     (c2_w - Cm(1.2)) / 2, img_h1, bg=RGBColor(0xCC, 0xCC, 0xBB),
                     image_url=beat.col2_image_url)
    _img_placeholder(slide, c2_x + Cm(0.5), img_y + img_h1 + Cm(0.15),
                     c2_w - Cm(1.0), img_h2, bg=RGBColor(0xCC, 0xCC, 0xBB),
                     image_url=beat.col2_image_url)

    # Col 3
    _add_text(slide, beat.mb_col3_title,
              c3_x + Cm(0.5), bar_y + Cm(0.4), c3_w - Cm(0.8), Cm(0.7),
              size=8, bold=True, color=NAVY, wrap=True, spacing=0.04)
    _add_text(slide, beat.mb_col3_accent,
              c3_x + Cm(0.5), bar_y + Cm(1.1), c3_w - Cm(0.8), Cm(0.5),
              size=8, bold=True, color=TEAL, spacing=0.02)
    _add_text(slide, beat.mb_col3_body,
              c3_x + Cm(0.5), bar_y + Cm(1.7), c3_w - Cm(0.8), Cm(3.5),
              size=7, color=GRAY_D, wrap=True, line_spacing=12)
    _img_placeholder(slide, c3_x + Cm(0.5), bar_y + Cm(5.4),
                     c3_w - Cm(1.0), Cm(3.4), bg=RGBColor(0xCC, 0xCC, 0xBB),
                     image_url=beat.col3_image_url)
    _img_placeholder(slide, c3_x + Cm(0.5), bar_y + Cm(9.0),
                     (c3_w - Cm(1.2)) / 2, Cm(1.8), bg=RGBColor(0xCC, 0xCC, 0xBB),
                     image_url=beat.col3_image_url)
    _img_placeholder(slide, c3_x + Cm(0.5) + (c3_w - Cm(1.2)) / 2 + Cm(0.15),
                     bar_y + Cm(9.0), (c3_w - Cm(1.2)) / 2, Cm(1.8),
                     bg=RGBColor(0xCC, 0xCC, 0xBB),
                     image_url=beat.col3_image_url)

    _page_num(slide, n, dark=True)


# ── Main builder ──────────────────────────────────────────────────────────

def build_indigo_pptx(story: IndigoStoryUnit) -> bytes:
    image_cache: dict[str, bytes | None] = {}
    jpeg_cache: dict[tuple[str, float], bytes | None] = {}
    raw_cache_token = _IMAGE_RAW_CACHE.set(image_cache)
    jpeg_cache_token = _IMAGE_JPEG_CACHE.set(jpeg_cache)
    try:
        _prefetch_story_images(story, image_cache)

        prs = Presentation()
        prs.slide_width = SW
        prs.slide_height = SH

        _slide01_cover(prs, story)
        _slide02_taglines(prs, story)
        _slide_cinematic(prs, 3, story,
                         BG_DARK,
                         story.taglines[0].zh + "  ·  " + story.taglines[0].sub,
                         story.concept_poem, "STORYLINE CONCEPT",
                         _beat_image(story, 1, "mood_image_url", "image_url", "col2_image_url", "col3_image_url"))
        for i in range(3):
            _slide_origin(prs, 4 + i, story, i)
        _slide_cinematic(prs, 7, story,
                         RGBColor(0x0E, 0x16, 0x10),
                         story.emotion_headline,
                         story.emotion_poem, "STORY EMOTION",
                         _beat_image(story, 3, "mood_image_url", "image_url", "col2_image_url", "col3_image_url"))
        _slide_story_summary(prs, 8, story)
        _slide_story_mapping(prs, 9, story)
        _slide_story_flow_grid(prs, 10, story)
        for i, beat in enumerate(story.beats):
            _slide_beat_cover(prs, 11 + i * 2, beat)
            _slide_moodboard(prs, 12 + i * 2, beat, story)

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
    finally:
        _IMAGE_JPEG_CACHE.reset(jpeg_cache_token)
        _IMAGE_RAW_CACHE.reset(raw_cache_token)
