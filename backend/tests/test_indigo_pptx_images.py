import base64
from collections import Counter
import io
import sys
import unittest
from pathlib import Path
from unittest.mock import Mock, patch

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from app.core.models import IndigoBeat, IndigoOrigin, IndigoStoryUnit, IndigoTagline  # noqa: E402
from app.services.indigo_pptx_builder import build_indigo_pptx  # noqa: E402


def _image_url(seed: int) -> str:
    color = ((seed * 45) % 255, (seed * 75) % 255, (seed * 105) % 255)
    image = Image.new("RGB", (1600, 900), color)
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    payload = base64.b64encode(buf.getvalue()).decode("ascii")
    return f"data:image/png;base64,{payload}"


def _beat(num: int) -> IndigoBeat:
    image = _image_url(num)
    return IndigoBeat(
        num=f"{num:02d}",
        name_zh=f"第{num}幕",
        space_zh=f"空间{num}",
        ghost_en=f"BEAT\n{num}",
        narrative=f"这是第{num}幕的故事叙述，用于检查导出图片是否被合理分配。",
        tagline=f"标语{num}",
        mb_ghost_en=f"MOOD\n{num}",
        mb_concept=f"概念{num}",
        mb_concept_sub=f"副标题{num}",
        mb_col2_title=f"第二栏标题{num}",
        mb_col2_accent=f"第二栏强调{num}",
        mb_col2_body=f"第二栏正文{num}",
        mb_col3_title=f"第三栏标题{num}",
        mb_col3_accent=f"第三栏强调{num}",
        mb_col3_body=f"第三栏正文{num}",
        image_url=image,
        mood_image_url=_image_url(num + 10),
        col2_image_url=_image_url(num + 20),
        col3_image_url=_image_url(num + 30),
    )


def _story() -> IndigoStoryUnit:
    return IndigoStoryUnit(
        city="上海",
        district="武康路",
        hotel_en="Shanghai Wukang Road",
        taglines=[
            IndigoTagline(zh="梧里光影", sub="街巷里的日常光线"),
            IndigoTagline(zh="路上旧梦", sub="老建筑与新生活"),
            IndigoTagline(zh="转角日常", sub="把街区带进酒店"),
        ],
        concept_poem=["第一段概念文本。", "第二段概念文本。"],
        origins=[
            IndigoOrigin(title="背景", headline="街区背景标题", body="街区背景正文。"),
            IndigoOrigin(title="人群", headline="人群生活标题", body="人群生活正文。"),
            IndigoOrigin(title="空间", headline="空间转译标题", body="空间转译正文。"),
        ],
        emotion_headline="「街」与「店」之间的日常流动",
        emotion_poem=["第一段情绪文本。", "第二段情绪文本。"],
        story_summary="以街区日常作为酒店触点的叙事线索。",
        beats=[_beat(i) for i in range(1, 7)],
    )


def _picture_count(slide) -> int:
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def _picture_frames(slide) -> list[tuple[int, int, int, int]]:
    return [
        (shape.left, shape.top, shape.width, shape.height)
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


class IndigoPptxImagesTest(unittest.TestCase):
    def test_first_ten_slides_allocate_available_beat_images(self) -> None:
        deck = Presentation(io.BytesIO(build_indigo_pptx(_story())))

        self.assertEqual(len(deck.slides), 22)
        for slide_no in range(1, 11):
            with self.subTest(slide_no=slide_no):
                self.assertGreater(_picture_count(deck.slides[slide_no - 1]), 0)

    def test_embedded_images_match_placeholder_aspect_ratio(self) -> None:
        deck = Presentation(io.BytesIO(build_indigo_pptx(_story())))

        for slide_no, slide in enumerate(deck.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                image_width, image_height = shape.image.size
                image_ratio = image_width / image_height
                frame_ratio = shape.width / shape.height
                with self.subTest(slide_no=slide_no, shape_id=shape.shape_id):
                    self.assertAlmostEqual(image_ratio, frame_ratio, delta=0.02)

    def test_origin_slides_use_varied_editorial_layouts(self) -> None:
        deck = Presentation(io.BytesIO(build_indigo_pptx(_story())))

        origin_frames = [_picture_frames(deck.slides[i])[0] for i in range(3, 6)]

        self.assertEqual(len(set(origin_frames)), 3)

    def test_story_touchpoints_use_flat_equal_grid(self) -> None:
        deck = Presentation(io.BytesIO(build_indigo_pptx(_story())))
        slide = deck.slides[9]
        frames = _picture_frames(slide)
        x_counts = Counter(frame[0] for frame in frames)
        y_counts = Counter(frame[1] for frame in frames)
        slide_text = "\n".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )

        self.assertEqual(len(frames), 6)
        self.assertEqual(sorted(x_counts.values()), [2, 2, 2])
        self.assertEqual(sorted(y_counts.values()), [3, 3])
        self.assertIn("空间触点", slide_text)
        self.assertNotIn("故事流线", slide_text)

    def test_repeated_remote_image_is_downloaded_once_per_deck(self) -> None:
        story = _story()
        image_url = "https://images.example.test/shared.png"
        for beat in story.beats:
            beat.image_url = image_url
            beat.mood_image_url = image_url
            beat.col2_image_url = image_url
            beat.col3_image_url = image_url

        image = Image.new("RGB", (1600, 900), (40, 80, 120))
        buf = io.BytesIO()
        image.save(buf, format="PNG")
        response = Mock(content=buf.getvalue())
        response.raise_for_status.return_value = None

        with patch("httpx.get", return_value=response) as get_image:
            deck = Presentation(io.BytesIO(build_indigo_pptx(story)))

        self.assertEqual(len(deck.slides), 22)
        self.assertEqual(get_image.call_count, 1)


if __name__ == "__main__":
    unittest.main()
